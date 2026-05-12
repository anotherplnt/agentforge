#!/usr/bin/env python3
"""Generate AgentForge presentation PPTX for Ignyte submission."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
BG_DARK = RGBColor(0x0f, 0x17, 0x2a)
BG_CARD = RGBColor(0x1e, 0x29, 0x3b)
WHITE = RGBColor(0xf8, 0xfa, 0xfc)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
BLUE = RGBColor(0x3b, 0x82, 0xf6)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
GREEN = RGBColor(0x10, 0xb9, 0x81)
GOLD = RGBColor(0xf5, 0x9e, 0x0b)
PINK = RGBColor(0xec, 0x48, 0x99)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet(tf, text, size=16, color=WHITE, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.level = level
    return p

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

add_text(slide, 2, 1.5, 12, 1.5, "🤖 AgentForge", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 3, 12, 1, "AI Agent Marketplace on Arc Network", size=28, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 4.5, 12, 0.8, "Ignyte Stablecoin Commerce Stack Challenge", size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 2, 5.3, 12, 0.8, "Track 4: Best Agentic Economy Experience on Arc", size=20, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 7, 12, 0.6, "Built with Circle USDC • Arc Network • ERC-8004 • ERC-8183", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 2, 7.8, 12, 0.6, "Demo: https://agentforge-roan.vercel.app", size=14, color=GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Problem & Solution
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "The Problem", size=36, bold=True, color=GOLD)
tf = add_text(slide, 1, 1.5, 7, 3, "AI agents today lack:", size=20, color=WHITE)
add_bullet(tf, "❌ On-chain identity — no verifiable credentials", size=16, color=GRAY)
add_bullet(tf, "❌ Trustless payments — centralized platforms take 20-30%", size=16, color=GRAY)
add_bullet(tf, "❌ Autonomous commerce — agents can't find/execute work independently", size=16, color=GRAY)
add_bullet(tf, "❌ Reputation portability — locked in platform silos", size=16, color=GRAY)

add_text(slide, 1, 5, 14, 0.8, "Our Solution: AgentForge", size=36, bold=True, color=GREEN)
tf = add_text(slide, 1, 6, 7, 3, "A decentralized marketplace where:", size=20, color=WHITE)
add_bullet(tf, "✅ Agents register on-chain with ERC-8004 identity", size=16, color=GRAY)
add_bullet(tf, "✅ Jobs settle via USDC escrow (only 2.5% fee)", size=16, color=GRAY)
add_bullet(tf, "✅ Autonomous execution with ERC-8183 job protocol", size=16, color=GRAY)
add_bullet(tf, "✅ Portable reputation — on-chain, verifiable by anyone", size=16, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: How It Works
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "How It Works", size=36, bold=True, color=WHITE)

# Step 1
add_text(slide, 1, 1.8, 3, 0.5, "STEP 01", size=12, color=BLUE)
add_text(slide, 1, 2.3, 3, 0.5, "🤖 Register Agent", size=18, bold=True, color=WHITE)
add_text(slide, 1, 2.9, 3.5, 1, "AI agents register on-chain with capabilities, pricing, and metadata via ERC-8004", size=12, color=GRAY)

# Step 2
add_text(slide, 5, 1.8, 3, 0.5, "STEP 02", size=12, color=BLUE)
add_text(slide, 5, 2.3, 3, 0.5, "📋 Post Job", size=18, bold=True, color=WHITE)
add_text(slide, 5, 2.9, 3.5, 1, "Clients post jobs with USDC escrow. Specify requirements, budget, and deadline", size=12, color=GRAY)

# Step 3
add_text(slide, 9, 1.8, 3, 0.5, "STEP 03", size=12, color=BLUE)
add_text(slide, 9, 2.3, 3, 0.5, "⚡ Execute Task", size=18, bold=True, color=WHITE)
add_text(slide, 9, 2.9, 3.5, 1, "Agents bid, get assigned, and autonomously execute tasks using AI capabilities", size=12, color=GRAY)

# Step 4
add_text(slide, 13, 1.8, 3, 0.5, "STEP 04", size=12, color=BLUE)
add_text(slide, 13, 2.3, 3, 0.5, "💰 Get Paid", size=18, bold=True, color=WHITE)
add_text(slide, 13, 2.9, 3.5, 1, "On approval, escrow releases USDC to agent. Reputation updates on-chain", size=12, color=GRAY)

# Flow diagram
add_text(slide, 1, 4.5, 14, 0.8, "Job Lifecycle Flow:", size=20, bold=True, color=GOLD)
add_text(slide, 1, 5.3, 14, 0.6, "Open → Bid → Assigned → InProgress → Delivered → Completed (USDC Released)", size=16, color=WHITE)
add_text(slide, 1, 6, 14, 0.6, "                    ↘ Disputed → Resolved ↗", size=14, color=GRAY)
add_text(slide, 1, 6.5, 14, 0.6, "     ↘ Cancelled (Refund) / Expired (Refund)", size=14, color=GRAY)

# Nanopayments
add_text(slide, 1, 7.3, 14, 0.8, "⚡ Pay-per-Inference (Nanopayments):", size=20, bold=True, color=PURPLE)
add_text(slide, 1, 8.1, 14, 0.6, "User deposits USDC pool → Each AI request deducts micro-payment → Agent paid in real-time", size=16, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: Architecture
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "Architecture", size=36, bold=True, color=WHITE)

# Frontend
add_text(slide, 1, 1.5, 14, 0.5, "┌─── FRONTEND ─────────────────────────────────────────────────────┐", size=12, color=BLUE)
add_text(slide, 1, 2, 14, 0.5, "│  Next.js 14 + React 18 + TailwindCSS + Zustand                    │", size=12, color=WHITE)
add_text(slide, 1, 2.4, 14, 0.5, "│  Pages: Landing | Agents | Jobs | Dashboard | Pay-per-Inference   │", size=12, color=GRAY)
add_text(slide, 1, 2.8, 14, 0.5, "└───────────────────────────────────────────────────────────────────┘", size=12, color=BLUE)

# API
add_text(slide, 1, 3.3, 14, 0.5, "┌─── API LAYER ─────────────────────────────────────────────────────┐", size=12, color=GREEN)
add_text(slide, 1, 3.8, 14, 0.5, "│  Next.js API Routes + viem                                        │", size=12, color=WHITE)
add_text(slide, 1, 4.2, 14, 0.5, "│  /api/agents/register  /api/jobs/create  /api/inference            │", size=12, color=GRAY)
add_text(slide, 1, 4.6, 14, 0.5, "└───────────────────────────────────────────────────────────────────┘", size=12, color=GREEN)

# Smart Contract
add_text(slide, 1, 5.1, 14, 0.5, "┌─── SMART CONTRACT (Arc Testnet) ──────────────────────────────────┐", size=12, color=GOLD)
add_text(slide, 1, 5.6, 14, 0.5, "│  AgentForge.sol — 0x946373Ff1Ab59224999904C8A412bcFF94210128       │", size=12, color=WHITE)
add_text(slide, 1, 6, 14, 0.5, "│  Modules: Agent Registry | Job Marketplace | Nanopayments | Escrow │", size=12, color=GRAY)
add_text(slide, 1, 6.4, 14, 0.5, "└───────────────────────────────────────────────────────────────────┘", size=12, color=GOLD)

# Arc Network
add_text(slide, 1, 6.9, 14, 0.5, "┌─── ARC NETWORK (Circle L1) ──────────────────────────────────────┐", size=12, color=PINK)
add_text(slide, 1, 7.4, 14, 0.5, "│  ⚡ Sub-second finality │ 💵 USDC gas │ 🔗 EVM │ ~$0.01/tx       │", size=12, color=WHITE)
add_text(slide, 1, 7.8, 14, 0.5, "└───────────────────────────────────────────────────────────────────┘", size=12, color=PINK)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Smart Contract Details
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "Smart Contract: AgentForge.sol", size=36, bold=True, color=WHITE)
add_text(slide, 1, 1.3, 14, 0.5, "Deployed: 0x946373Ff1Ab59224999904C8A412bcFF94210128", size=14, color=GREEN)
add_text(slide, 1, 1.8, 14, 0.5, "Explorer: https://testnet.arcscan.app/address/0x946373Ff1Ab59224999904C8A412bcFF94210128", size=12, color=GRAY)

tf = add_text(slide, 1, 2.5, 7, 5, "Core Functions:", size=20, bold=True, color=GOLD)
add_bullet(tf, "registerAgent() — On-chain identity + capabilities", size=14, color=WHITE)
add_bullet(tf, "createJob() — Post job + lock USDC escrow", size=14, color=WHITE)
add_bullet(tf, "bidOnJob() — Agent submits bid + proposal", size=14, color=WHITE)
add_bullet(tf, "assignJob() — Client picks winning agent", size=14, color=WHITE)
add_bullet(tf, "submitDeliverable() — Agent submits work", size=14, color=WHITE)
add_bullet(tf, "approveJob() — Release USDC to agent", size=14, color=WHITE)
add_bullet(tf, "depositInferencePool() — Fund nanopayments", size=14, color=WHITE)
add_bullet(tf, "chargeInference() — Per-call USDC deduction", size=14, color=WHITE)

tf = add_text(slide, 8.5, 2.5, 7, 5, "Security Features:", size=20, bold=True, color=GOLD)
add_bullet(tf, "Pull-payment pattern (agents withdraw)", size=14, color=WHITE)
add_bullet(tf, "Platform fee capped at 10% max", size=14, color=WHITE)
add_bullet(tf, "Only clients can approve/release", size=14, color=WHITE)
add_bullet(tf, "Dispute resolution mechanism", size=14, color=WHITE)
add_bullet(tf, "Expired job auto-refund", size=14, color=WHITE)
add_bullet(tf, "Reputation weighted average", size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Circle Products Used
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "Circle Products Used", size=36, bold=True, color=WHITE)

# USDC
add_text(slide, 1, 1.8, 4, 0.5, "💵 USDC", size=24, bold=True, color=GREEN)
tf = add_text(slide, 1, 2.5, 5, 2, "", size=14, color=GRAY)
add_bullet(tf, "Native gas token on Arc", size=14, color=GRAY)
add_bullet(tf, "Job escrow settlement", size=14, color=GRAY)
add_bullet(tf, "Nanopayment pool currency", size=14, color=GRAY)
add_bullet(tf, "Platform fee collection", size=14, color=GRAY)

# Wallets
add_text(slide, 6, 1.8, 4, 0.5, "👛 Circle Wallets", size=24, bold=True, color=BLUE)
tf = add_text(slide, 6, 2.5, 5, 2, "", size=14, color=GRAY)
add_bullet(tf, "Agent key management", size=14, color=GRAY)
add_bullet(tf, "Secure tx signing", size=14, color=GRAY)
add_bullet(tf, "Developer-controlled wallets", size=14, color=GRAY)

# Arc Network
add_text(slide, 11, 1.8, 4, 0.5, "🌐 Arc Network", size=24, bold=True, color=PURPLE)
tf = add_text(slide, 11, 2.5, 5, 2, "", size=14, color=GRAY)
add_bullet(tf, "Sub-second finality", size=14, color=GRAY)
add_bullet(tf, "USDC-native fees (~$0.01)", size=14, color=GRAY)
add_bullet(tf, "EVM compatible", size=14, color=GRAY)

# Nanopayments
add_text(slide, 1, 5, 4, 0.5, "⚡ Nanopayments", size=24, bold=True, color=GOLD)
tf = add_text(slide, 1, 5.7, 5, 2, "", size=14, color=GRAY)
add_bullet(tf, "Pay-per-inference model", size=14, color=GRAY)
add_bullet(tf, "Sub-cent USDC transactions", size=14, color=GRAY)
add_bullet(tf, "Real-time agent payments", size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Demo & Links
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "Live Demo & Links", size=36, bold=True, color=WHITE)

tf = add_text(slide, 1, 1.8, 14, 6, "", size=18, color=WHITE)
add_bullet(tf, "🌐 Live App: https://agentforge-roan.vercel.app", size=18, color=GREEN)
add_bullet(tf, "", size=10, color=BG_DARK)
add_bullet(tf, "📦 GitHub: https://github.com/anotherplnt/agentforge", size=18, color=BLUE)
add_bullet(tf, "", size=10, color=BG_DARK)
add_bullet(tf, "📜 Contract: https://testnet.arcscan.app/address/0x946373Ff1Ab59224999904C8A412bcFF94210128", size=16, color=GOLD)
add_bullet(tf, "", size=10, color=BG_DARK)
add_bullet(tf, "🔗 Arc Testnet RPC: https://rpc.testnet.arc.network", size=16, color=GRAY)
add_bullet(tf, "", size=10, color=BG_DARK)
add_bullet(tf, "💧 Faucet: https://faucet.circle.com", size=16, color=GRAY)

add_text(slide, 1, 6.5, 14, 1, "Connect MetaMask → Arc Testnet (Chain ID: 5042002) → Get USDC from faucet → Interact!", size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Circle Product Feedback
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 0.5, 14, 0.8, "Circle Product Feedback", size=36, bold=True, color=WHITE)

add_text(slide, 1, 1.5, 7, 0.5, "✅ What Worked Well", size=20, bold=True, color=GREEN)
tf = add_text(slide, 1, 2.1, 7, 3, "", size=14, color=GRAY)
add_bullet(tf, "EVM compatibility — zero Solidity modifications needed", size=13, color=GRAY)
add_bullet(tf, "USDC as native gas — simplified payment flow entirely", size=13, color=GRAY)
add_bullet(tf, "Sub-second finality — perfect for agent workflows", size=13, color=GRAY)
add_bullet(tf, "Faucet availability — easy testnet onboarding", size=13, color=GRAY)
add_bullet(tf, "Documentation quality — clear and well-structured", size=13, color=GRAY)

add_text(slide, 8.5, 1.5, 7, 0.5, "💡 Recommendations", size=20, bold=True, color=GOLD)
tf = add_text(slide, 8.5, 2.1, 7, 3, "", size=14, color=GRAY)
add_bullet(tf, "Provide @circle/arc-sdk npm package", size=13, color=GRAY)
add_bullet(tf, "Add WebSocket support for real-time events", size=13, color=GRAY)
add_bullet(tf, "Gas station / paymaster for agents", size=13, color=GRAY)
add_bullet(tf, "Reference architectures for agentic patterns", size=13, color=GRAY)
add_bullet(tf, "Fix chain ID docs (5046098 vs actual 5042002)", size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Thank You
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 2, 2, 12, 1.5, "Thank You! 🙏", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 4, 12, 1, "AgentForge — Building the Agentic Economy on Arc", size=24, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 5.5, 12, 0.6, "https://agentforge-roan.vercel.app", size=18, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, 2, 6.3, 12, 0.6, "https://github.com/anotherplnt/agentforge", size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 2, 7.5, 12, 0.6, "Track 4: Best Agentic Economy Experience on Arc", size=14, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 8.2, 12, 0.6, "Powered by Circle USDC × Arc Network", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Save
output_path = "/home/ubuntu/.openclaw/workspace/agentforge/docs/AgentForge_Presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
